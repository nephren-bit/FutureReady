"""
Unit tests for Layer 1 extractors. Each test builds a minimal, valid fixture
file on the fly (rather than shipping binary test assets) using the same
libraries the extractors themselves depend on. Tests for optional/heavy
dependencies (OpenCV) are skipped automatically if the dependency isn't
installed, so the suite still runs on a minimal install.
"""

from __future__ import annotations

from pathlib import Path

import numpy as np
import pytest

from extractors.audio_extractor import AudioExtractor
from extractors.pdf_extractor import PDFExtractor
from extractors.ppt_extractor import PPTExtractor
from models.features import AudioFeature, ResumeFeature, SlideFeature, VideoFeature


@pytest.fixture()
def sample_pdf_path(tmp_path: Path) -> Path:
    fitz = pytest.importorskip("fitz")
    pdf_path = tmp_path / "resume.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "John Doe", fontsize=18)
    page.insert_text((72, 110), "Skills", fontsize=13)
    page.insert_text((72, 130), "Python, SQL, Leadership", fontsize=10)
    page.insert_text((72, 160), "Experience", fontsize=13)
    page.insert_text((72, 180), "Led a team of engineers to ship a product.", fontsize=10)
    document.save(str(pdf_path))
    document.close()
    return pdf_path


@pytest.fixture()
def sample_pptx_path(tmp_path: Path) -> Path:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Introduction"
    body = slide.placeholders[1]
    body.text_frame.text = "Welcome to the presentation."

    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture()
def sample_wav_path(tmp_path: Path) -> Path:
    soundfile = pytest.importorskip("soundfile")
    sample_rate = 22050
    duration_sec = 2.0
    t = np.linspace(0, duration_sec, int(sample_rate * duration_sec), endpoint=False)
    tone = 0.2 * np.sin(2 * np.pi * 220.0 * t)
    wav_path = tmp_path / "speech.wav"
    soundfile.write(str(wav_path), tone.astype(np.float32), sample_rate)
    return wav_path


class TestPDFExtractor:
    def test_extract_returns_resume_feature(self, sample_pdf_path: Path) -> None:
        result = PDFExtractor().extract(sample_pdf_path)
        assert isinstance(result, ResumeFeature)
        assert result.page_count == 1
        assert "John Doe" in result.text
        assert result.word_count > 0

    def test_extract_raises_on_invalid_file(self, tmp_path: Path) -> None:
        pytest.importorskip("fitz")
        bad_path = tmp_path / "not_a_pdf.pdf"
        bad_path.write_text("this is not a real pdf")
        with pytest.raises(RuntimeError):
            PDFExtractor().extract(bad_path)


class TestPPTExtractor:
    def test_extract_returns_slide_feature(self, sample_pptx_path: Path) -> None:
        result = PPTExtractor().extract(sample_pptx_path)
        assert isinstance(result, SlideFeature)
        assert result.slide_count == 1
        assert result.slides[0].title == "Introduction"


class TestAudioExtractor:
    def test_extract_returns_audio_feature(self, sample_wav_path: Path) -> None:
        result = AudioExtractor().extract(sample_wav_path)
        assert isinstance(result, AudioFeature)
        assert result.sample_rate == 22050
        assert result.duration_sec == pytest.approx(2.0, abs=0.1)

    def test_extract_raises_on_invalid_file(self, tmp_path: Path) -> None:
        bad_path = tmp_path / "not_audio.wav"
        bad_path.write_bytes(b"not a real wav file")
        with pytest.raises(RuntimeError):
            AudioExtractor().extract(bad_path)


class TestVideoExtractor:
    def test_extract_returns_video_feature(self, tmp_path: Path) -> None:
        cv2 = pytest.importorskip("cv2")
        from extractors.video_extractor import VideoExtractor

        video_path = tmp_path / "clip.mp4"
        writer = cv2.VideoWriter(
            str(video_path), cv2.VideoWriter_fourcc(*"mp4v"), 10.0, (64, 64)
        )
        for i in range(20):
            frame = np.full((64, 64, 3), fill_value=(i * 10) % 256, dtype=np.uint8)
            writer.write(frame)
        writer.release()

        if not video_path.exists() or video_path.stat().st_size == 0:
            pytest.skip("OpenCV could not encode an MP4 in this environment (missing codec).")

        result = VideoExtractor(sample_count=5).extract(video_path)
        assert isinstance(result, VideoFeature)
        assert result.frame_count == 20
        assert result.sampled_frame_count > 0
