"""
extractors/ppt_extractor.py

Extracts structured information from a presentation (.pptx) file using
python-pptx. Pulls out per-slide titles, bullet text, image/chart/table
counts, font and color usage, speaker notes, and average text length. The
result is a typed `SlideFeature` model — Layer 1 only, no AI reasoning.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from extractors.base import BaseExtractor
from models.features import SlideFeature, SlideInfo
from utils.logger import get_logger

logger = get_logger(__name__)


class PPTExtractor(BaseExtractor[SlideFeature]):
    """Extracts structured data from a PPTX presentation file (Layer 1)."""

    def extract(self, file_path: Path) -> SlideFeature:
        """
        Run full extraction on a PPTX file.

        Args:
            file_path: Path to the .pptx file on disk.

        Returns:
            A `SlideFeature` model describing the presentation.

        Raises:
            RuntimeError: If the file cannot be opened or parsed.
        """
        try:
            presentation = Presentation(str(file_path))
        except Exception as exc:  # noqa: BLE001
            logger.exception("Failed to open PPTX %s", file_path)
            raise RuntimeError(f"Could not open PPTX file: {exc}") from exc

        slides_data: list[SlideInfo] = []
        total_image_count = 0
        total_chart_count = 0
        total_table_count = 0
        fonts_used: set[str] = set()
        colors_used: set[str] = set()
        text_lengths: list[int] = []

        for index, slide in enumerate(presentation.slides):
            slide_info = self._extract_slide(slide, index + 1)
            slides_data.append(slide_info)

            total_image_count += slide_info.image_count
            total_chart_count += slide_info.chart_count
            total_table_count += slide_info.table_count
            fonts_used.update(slide_info.fonts)
            colors_used.update(slide_info.colors)
            text_lengths.append(slide_info.text_length)

        avg_text_length = (
            round(sum(text_lengths) / len(text_lengths), 2) if text_lengths else 0.0
        )

        return SlideFeature(
            slide_count=len(presentation.slides),
            slides=slides_data,
            image_count=total_image_count,
            chart_count=total_chart_count,
            table_count=total_table_count,
            fonts=sorted(fonts_used),
            colors=sorted(colors_used),
            average_text_length=avg_text_length,
        )

    def _extract_slide(self, slide: Any, slide_number: int) -> SlideInfo:
        """Extract data for a single slide."""
        title = ""
        bullets: list[str] = []
        notes = ""
        image_count = 0
        chart_count = 0
        table_count = 0
        fonts: set[str] = set()
        colors: set[str] = set()
        text_chars = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if shape.has_chart:
                chart_count += 1
            if shape.has_table:
                table_count += 1

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join(run.text for run in paragraph.runs)
                    if not paragraph_text.strip():
                        continue

                    text_chars += len(paragraph_text)

                    if shape == slide.shapes.title:
                        title = paragraph_text
                    else:
                        bullets.append(paragraph_text)

                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        try:
                            if run.font.color and run.font.color.type is not None:
                                colors.add(str(run.font.color.rgb))
                        except (AttributeError, TypeError):
                            # Color may be theme-based or unset; skip gracefully.
                            pass

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        return SlideInfo(
            slide_number=slide_number,
            title=title,
            bullets=bullets,
            notes=notes,
            image_count=image_count,
            chart_count=chart_count,
            table_count=table_count,
            fonts=sorted(fonts),
            colors=sorted(colors),
            text_length=text_chars,
        )
